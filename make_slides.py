"""
Laoban: 5/10/K engine — VISUAL deck.
Philosophy: slides show, the speaker tells. Near-zero prose on slides;
all narration lives in speaker notes. Real rendered image assets
(gradient felt, glows, matplotlib charts/heatmap) + crisp vector card art.
"""
import os, math, random
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.colors import ListedColormap
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

random.seed(7); np.random.seed(7)
ASSETS = "assets"; os.makedirs(ASSETS, exist_ok=True)

# ── Palette ──────────────────────────────────────────────────────────────────
DARK   = RGBColor(0x12, 0x18, 0x2B)
MID    = RGBColor(0x1E, 0x2B, 0x47)
PANEL  = RGBColor(0x29, 0x3A, 0x5E)
GOLD   = RGBColor(0xF5, 0xC5, 0x18)
RED    = RGBColor(0xD7, 0x26, 0x3D)
GREEN  = RGBColor(0x1E, 0x7A, 0x46)
TEAL   = RGBColor(0x27, 0xC4, 0xA6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHT = RGBColor(0xEC, 0xEF, 0xF4)
INK    = RGBColor(0x14, 0x14, 0x16)
MUTED  = RGBColor(0x8E, 0x9F, 0xBC)
GREY   = RGBColor(0xCF, 0xD6, 0xE0)
HEX = dict(gold="#F5C518", red="#D7263D", teal="#27C4A6", ink="#141416",
           off="#ECEFF4", muted="#8E9FBC", panel="#293A5E", greyl="#5A6B86")

# ════════════════════════════════════════════════════════════════════════════
#  ASSET GENERATION
# ════════════════════════════════════════════════════════════════════════════
def gradient_bg(path, center_rgb, edge_rgb, vignette=0.0, w=2000, h=1125):
    yy, xx = np.mgrid[0:h, 0:w].astype(float)
    cx, cy = w/2, h/2
    r = np.sqrt(((xx-cx)/(w/2))**2 + ((yy-cy)/(h/2))**2)
    r = np.clip(r/1.15, 0, 1)
    img = np.zeros((h, w, 3), float)
    for k in range(3):
        img[..., k] = center_rgb[k]*(1-r) + edge_rgb[k]*r
    if vignette:
        img *= (1 - vignette*np.clip(r-0.35, 0, 1)**1.6)[..., None]
    Image.fromarray(np.clip(img, 0, 255).astype("uint8")).save(path)

def glow(path, rgb, size=900):
    yy, xx = np.mgrid[0:size, 0:size].astype(float)
    cx = cy = size/2
    r = np.sqrt(((xx-cx)/(size/2))**2 + ((yy-cy)/(size/2))**2)
    a = np.clip(np.exp(-(r**2)*3.2), 0, 1)*0.85
    rgba = np.zeros((size, size, 4), float)
    rgba[..., 0], rgba[..., 1], rgba[..., 2] = rgb
    rgba[..., 3] = a*255
    Image.fromarray(rgba.astype("uint8")).save(path)

def style_dark(ax):
    ax.set_facecolor("none")
    for sp in ax.spines.values(): sp.set_visible(False)
    ax.tick_params(colors=HEX["muted"], labelsize=13)

def chart_score():
    fig, ax = plt.subplots(figsize=(7.6, 5.0), dpi=200)
    fig.patch.set_alpha(0)
    bars = ax.bar([0, 1], [60, 400], width=0.46,
                  color=[HEX["teal"], HEX["red"]], zorder=3)
    for x, v, c in [(0, 60, HEX["teal"]), (1, 400, HEX["red"])]:
        ax.text(x, v+14, str(v), ha="center", va="bottom", color=c,
                fontsize=46, fontweight="bold", family="Georgia")
    ax.text(0, -34, "WINS",  ha="center", color=HEX["muted"], fontsize=17, fontweight="bold")
    ax.text(1, -34, "LOSSES", ha="center", color=HEX["muted"], fontsize=17, fontweight="bold")
    ax.set_ylim(0, 470); ax.set_xlim(-0.6, 1.6)
    ax.axis("off")
    fig.savefig(f"{ASSETS}/chart_score.png", transparent=True, bbox_inches="tight")
    plt.close(fig)

def chart_dmc_ppo():
    fig, (a1, a2) = plt.subplots(1, 2, figsize=(11.2, 4.5), dpi=200)
    fig.patch.set_alpha(0)
    x = np.linspace(0, 1, 220)
    # DMC: high variance, barely improving
    mean1 = 0.12*x
    band = 0.42
    noisy = mean1 + np.random.uniform(-0.4, 0.4, x.size)
    a1.fill_between(x, mean1-band, mean1+band, color=HEX["red"], alpha=0.15)
    a1.plot(x, noisy, color=HEX["red"], lw=1.4)
    a1.set_title("DMC — high variance", color=HEX["red"], fontsize=18,
                 fontweight="bold", family="Georgia", pad=12)
    # PPO: smooth converging
    mean2 = 0.92*(1-np.exp(-3.6*x))
    a2.fill_between(x, mean2-0.05, mean2+0.05, color=HEX["teal"], alpha=0.20)
    a2.plot(x, mean2 + np.random.normal(0, 0.012, x.size), color=HEX["teal"], lw=2.6)
    a2.set_title("PPO — stable, converges", color=HEX["teal"], fontsize=18,
                 fontweight="bold", family="Georgia", pad=12)
    for ax in (a1, a2):
        style_dark(ax)
        ax.set_ylim(-0.55, 1.05); ax.set_xlim(0, 1)
        ax.set_xticks([]); ax.set_yticks([])
        ax.axhline(0, color=HEX["panel"], lw=1)
        ax.set_xlabel("training  →", color=HEX["muted"], fontsize=12)
    a1.set_ylabel("reward signal", color=HEX["muted"], fontsize=12)
    fig.savefig(f"{ASSETS}/chart_dmc.png", transparent=True, bbox_inches="tight")
    plt.close(fig)

def chart_bench():
    fig, ax = plt.subplots(figsize=(9.4, 4.6), dpi=200)
    fig.patch.set_alpha(0)
    xs = [0, 1, 2]; heights = [98, 90, 50]
    cols = [HEX["teal"], HEX["gold"], HEX["red"]]
    bars = ax.bar(xs, heights, width=0.55, color=cols, zorder=3)
    bars[1].set_alpha(0.45); bars[1].set_hatch("////"); bars[1].set_edgecolor(HEX["gold"])
    labels_top = ["98%", "almost\nalways", "50 / 50"]
    for x, v, t, c in zip(xs, heights, labels_top, cols):
        ax.text(x, v+3, t, ha="center", va="bottom", color=c,
                fontsize=22, fontweight="bold", family="Georgia")
    cats = ["vs Random", "vs Novice /\nIntermediate", "vs Experienced"]
    for x, t in zip(xs, cats):
        ax.text(x, -8, t, ha="center", va="top", color=HEX["ink"], fontsize=15)
    ax.text(1, 44, "(not formally\nbenchmarked)", ha="center", color=HEX["greyl"],
            fontsize=10, style="italic")
    ax.set_ylim(0, 118); ax.set_xlim(-0.6, 2.6); ax.axis("off")
    fig.savefig(f"{ASSETS}/chart_bench.png", transparent=True, bbox_inches="tight")
    plt.close(fig)

def heat_features():
    cols_n, rows_n, total = 22, 8, 172
    C = np.full((rows_n, cols_n), np.nan)
    for idx in range(total):
        C[idx // cols_n, idx % cols_n] = 0 if idx < 60 else 1
    C = np.ma.masked_invalid(C)
    cmap = ListedColormap([HEX["gold"], HEX["teal"]])
    fig, ax = plt.subplots(figsize=(11.0, 4.2), dpi=200)
    fig.patch.set_alpha(0)
    ax.pcolormesh(C, cmap=cmap, vmin=0, vmax=1,
                  edgecolors="white", linewidth=2.0)
    ax.set_aspect("equal"); ax.invert_yaxis(); ax.axis("off")
    fig.savefig(f"{ASSETS}/heat_features.png", transparent=True, bbox_inches="tight")
    plt.close(fig)

# build assets
gradient_bg(f"{ASSETS}/bg_dark.png",  (0x22, 0x30, 0x50), (0x0E, 0x13, 0x22), vignette=0.45)
gradient_bg(f"{ASSETS}/bg_light.png", (0xFF, 0xFF, 0xFF), (0xE7, 0xEC, 0xF3), vignette=0.0)
glow(f"{ASSETS}/glow_gold.png", (0xF5, 0xC5, 0x18))
glow(f"{ASSETS}/glow_teal.png", (0x27, 0xC4, 0xA6))
glow(f"{ASSETS}/glow_red.png",  (0xD7, 0x26, 0x3D))
chart_score(); chart_dmc_ppo(); chart_bench(); heat_features()

# ════════════════════════════════════════════════════════════════════════════
#  DECK
# ════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
W, H = 13.33, 7.5
def i(v): return Inches(v)

def slide(dark=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_picture(f"{ASSETS}/{'bg_dark' if dark else 'bg_light'}.png",
                         0, 0, width=i(W), height=i(H))
    return s

def notes(s, t): s.notes_slide.notes_text_frame.text = t

def pic(s, path, x, y, w):
    return s.shapes.add_picture(path, i(x), i(y), width=i(w))

def pic_center(s, path, cy, w, x=None):
    p = s.shapes.add_picture(path, i(0), i(0), width=i(w))
    ph = p.height / 914400
    p.left = i((W-w)/2 if x is None else x); p.top = i(cy-ph/2)
    return p

def glowbehind(s, kind, cx, cy, d):
    s.shapes.add_picture(f"{ASSETS}/glow_{kind}.png", i(cx-d/2), i(cy-d/2),
                         width=i(d), height=i(d))

def rect(s, x, y, w, h, color, line=None, lw=1.0, rounded=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             i(x), i(y), i(w), i(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False; return shp

def oval(s, x, y, w, h, color, line=None, lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, i(x), i(y), i(w), i(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False; return shp

def txt(s, text, x, y, w, h, size, color, bold=False, italic=False,
        align=PP_ALIGN.LEFT, font="Trebuchet MS", anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(i(x), i(y), i(w), i(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for k, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return tb

def kicker(s, text, dark=True):
    """tiny eyebrow label, top-left — the only consistent on-slide chrome"""
    txt(s, text, 0.9, 0.6, 11, 0.5, 15, GOLD, bold=True, font="Trebuchet MS")

def arrow(s, x1, y1, x2, y2, color=GOLD, lw=2.5):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, i(x1), i(y1), i(x2), i(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(lw)
    ln = cn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    cn.shadow.inherit = False; return cn

# ── card art ─────────────────────────────────────────────────────────────────
def card(s, x, y, w=1.25, h=1.82, rank="A", suit="♠", point=False, facedown=False):
    if facedown:
        rect(s, x, y, w, h, MID, line=GOLD, lw=1.5, rounded=True)
        rect(s, x+0.12, y+0.12, w-0.24, h-0.24, DARK, line=GOLD, lw=0.75, rounded=True)
        txt(s, "♦", x, y, w, h, w*26, GOLD, align=PP_ALIGN.CENTER,
            font="Arial", anchor=MSO_ANCHOR.MIDDLE)
        return
    rect(s, x, y, w, h, WHITE, line=(GOLD if point else GREY),
         lw=(2.75 if point else 1.0), rounded=True)
    col = RED if suit in ("♥", "♦") else INK
    rs = max(9, w*13)
    txt(s, rank, x+0.08, y+0.05, w*0.7, 0.32, rs, col, bold=True, font="Georgia")
    txt(s, suit, x, y, w, h, w*34, col, align=PP_ALIGN.CENTER, font="Arial", anchor=MSO_ANCHOR.MIDDLE)
    if point:
        oval(s, x+w-0.27, y+0.09, 0.17, 0.17, GOLD)

def joker(s, x, y, w=1.25, h=1.82, red=True):
    rect(s, x, y, w, h, WHITE, line=GOLD, lw=2.0, rounded=True)
    col = RED if red else INK
    txt(s, "★", x, y+0.16, w, 0.5, w*26, col, align=PP_ALIGN.CENTER, font="Arial")
    txt(s, "JOKER", x, y+h*0.54, w, 0.4, w*8.5, col, bold=True, align=PP_ALIGN.CENTER, font="Georgia")


# ════════════════════════════════════════════════════════════════════════════
# 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s = slide(dark=True)
glowbehind(s, "gold", W/2, 1.7, 7)
txt(s, "LAOBAN", 0, 0.85, W, 1.7, 120, GOLD, bold=True, align=PP_ALIGN.CENTER, font="Georgia")
txt(s, "老板", 0, 2.55, W, 0.7, 26, OFFWHT, align=PP_ALIGN.CENTER, font="Georgia")
hand = [("3","♣",False),("7","♦",False),("10","♥",True),("K","♠",True),("A","♥",False),("2","♣",False)]
n = len(hand); cx = W/2; cw, ch = 1.45, 2.1; sp = 1.35; lift = 0.16
for k,(rk,su,pt) in enumerate(hand):
    t = k-(n-1)/2
    card(s, cx+t*sp-cw/2, 4.3+(t*t)*lift, cw, ch, rk, su, point=pt)
notes(s, "Hi, I'm Gavin. Today I want to talk about 5/10/K — wushiK — a card game with tens of "
         "millions of players worldwide that almost nobody in America knows. Introducing Laoban: "
         "the world's first engine designed to play it.")

# ════════════════════════════════════════════════════════════════════════════
# 2 — POPULARITY BLIND SPOT  (two glowing stats)
# ════════════════════════════════════════════════════════════════════════════
s = slide(dark=True)
kicker(s, "THE BLIND SPOT")
glowbehind(s, "teal", 3.6, 3.6, 5.5)
txt(s, "10M+", 0.6, 2.5, 6, 1.7, 110, TEAL, bold=True, align=PP_ALIGN.CENTER, font="Georgia")
txt(s, "players worldwide", 0.6, 4.3, 6, 0.5, 18, MUTED, align=PP_ALIGN.CENTER)
txt(s, "vs", 6.4, 3.3, 0.6, 0.6, 28, GREY, align=PP_ALIGN.CENTER, italic=True)
glowbehind(s, "red", 9.7, 3.6, 4.5)
txt(s, "≈ 0", 7.3, 2.5, 6, 1.7, 110, RED, bold=True, align=PP_ALIGN.CENTER, font="Georgia")
txt(s, "engines that exist for it", 7.3, 4.3, 6, 0.5, 18, MUTED, align=PP_ALIGN.CENTER)
notes(s, "What makes it fun: a large action space with imperfect information — you weigh taking points "
         "now versus saving them, while reading your opponent like in poker. And yet, despite all those "
         "players, there was no engine. That was the bottleneck, and honestly it surprised me.")

# ════════════════════════════════════════════════════════════════════════════
# 3 — THE SCORE  (chart image)
# ════════════════════════════════════════════════════════════════════════════
s = slide(dark=True)
kicker(s, "ME vs MY ROOMMATES, FALL QUARTER")
pic_center(s, f"{ASSETS}/chart_score.png", cy=4.05, w=7.6)
txt(s, '"the only way you\'d beat us is with an RL policy"', 0, 6.6, W, 0.5, 16,
    GOLD, italic=True, align=PP_ALIGN.CENTER, font="Georgia")
notes(s, "And I'm bad. By the end of fall quarter I'd won about 60 games and lost about 400. One of my "
         "roommates told me the only way I'd beat them is with an RL policy on my side. He may regret that.")

# ════════════════════════════════════════════════════════════════════════════
# 4 — WHY IT'S HARD  (your hand → branches | hidden opponent)
# ════════════════════════════════════════════════════════════════════════════
s = slide(dark=True)
kicker(s, "YOUR TURN")
for k,(rk,su,pt) in enumerate([("9","♠",False),("9","♥",False),("5","♦",True),("K","♣",True)]):
    card(s, 1.0+k*0.6, 2.9, 1.05, 1.5, rk, su, point=pt)
txt(s, "your hand", 1.0, 4.5, 3, 0.4, 14, MUTED, italic=True)
for k,o in enumerate(["play single","play pair","save the 5","drop the K","bomb?"]):
    yy = 2.25+k*0.62
    arrow(s, 3.7, 3.65, 5.3, yy+0.2, color=PANEL, lw=1.4)
    rect(s, 5.35, yy, 2.1, 0.46, MID, rounded=True)
    txt(s, o, 5.35, yy, 2.1, 0.46, 13, OFFWHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for k in range(4):
    card(s, 9.5+k*0.5, 2.9, 1.05, 1.5, "", "", facedown=True)
txt(s, "what do they have?", 9.0, 4.5, 4, 0.4, 14, RED, italic=True, align=PP_ALIGN.CENTER)
notes(s, "Every hand is a fork: dozens of legal plays, and you never see the opponent's cards. So each "
         "move is also a read — bluff, hold points, or commit. Big action space, imperfect information.")

# ════════════════════════════════════════════════════════════════════════════
# 5 — THE GAP  (peers have engines, this one doesn't)
# ════════════════════════════════════════════════════════════════════════════
s = slide(dark=False)
kicker(s, "EVERY NEIGHBOR HAS ONE")
peers = [("DouDiZhu","DouZero",True),("Sheng Ji","papers",True),
         ("Poker","Libratus",True),("5/10/K","—",False)]
cw2=2.7; gap=0.45; tot=len(peers)*cw2+(len(peers)-1)*gap; x0=(W-tot)/2
for k,(name,tag,has) in enumerate(peers):
    x = x0+k*(cw2+gap)
    rect(s, x, 2.0, cw2, 3.6, DARK if has else RED, rounded=True)
    txt(s, name, x, 2.35, cw2, 0.7, 21, WHITE, bold=True, align=PP_ALIGN.CENTER, font="Georgia")
    oval(s, x+cw2/2-0.5, 3.3, 1.0, 1.0, TEAL if has else WHITE)
    txt(s, "✓" if has else "✗", x+cw2/2-0.5, 3.3, 1.0, 1.0, 40, DARK if has else RED,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, tag, x, 4.65, cw2, 0.5, 14, GREY if has else WHITE, italic=True, align=PP_ALIGN.CENTER)
notes(s, "Adjacent games are well served. DouDiZhu has DouZero; other Sheng Ji variants have papers; "
         "poker has Libratus and others. 5/10/K had nothing — so I started there.")

# ════════════════════════════════════════════════════════════════════════════
# 6 — DMC → PPO  (chart image)
# ════════════════════════════════════════════════════════════════════════════
s = slide(dark=True)
kicker(s, "v1  →  v2")
pic_center(s, f"{ASSETS}/chart_dmc.png", cy=4.0, w=11.6)
notes(s, "My first approach was Deep Monte Carlo, like DouZero — zero bias, but it tries to model the "
         "entire outcome of each action. For a move set this large that's extremely noisy and high-variance. "
         "So I switched to a PPO-based policy: far more stable.")

# ════════════════════════════════════════════════════════════════════════════
# 7 — GARBAGE IN / GARBAGE OUT  (pipeline)
# ════════════════════════════════════════════════════════════════════════════
s = slide(dark=True)
kicker(s, "THE REAL BOTTLENECK")
def node(x, title, sub):
    rect(s, x, 3.0, 3.4, 2.0, MID, rounded=True)
    txt(s, title, x, 3.3, 3.4, 0.6, 22, OFFWHT, bold=True, align=PP_ALIGN.CENTER, font="Georgia")
    txt(s, sub, x, 3.95, 3.4, 0.7, 13, MUTED, align=PP_ALIGN.CENTER)
node(0.95, "opponent pool", "the training data")
node(5.0, "PPO policy", "learns from play")
node(9.05, "the bot", "only as good as\nits sparring")
arrow(s, 4.4, 4.0, 4.95, 4.0); arrow(s, 8.45, 4.0, 9.0, 4.0)
glowbehind(s, "gold", W/2, 6.0, 6)
txt(s, "garbage in  →  garbage out", 0, 5.6, W, 0.7, 26, GOLD, bold=True, italic=True,
    align=PP_ALIGN.CENTER, font="Georgia")
notes(s, "The hardest part wasn't the model — it was the data, which here means the opponents. A policy is "
         "only as good as what it trains against. Garbage in, garbage out. This is where I spent the most time.")

# ════════════════════════════════════════════════════════════════════════════
# 8 — THE DOOM LOOP  (cyclic degradation)
# ════════════════════════════════════════════════════════════════════════════
s = slide(dark=False)
kicker(s, "WHY NAIVE SELF-PLAY BROKE")
cxc, cyc, R = 6.66, 4.0, 1.95
pts = []
for a in (-90, 30, 150):
    pts.append((cxc+R*math.cos(math.radians(a)), cyc+R*math.sin(math.radians(a))))
for (nx,ny),lab,c in zip(pts, ["Strategy A","Strategy B","Strategy C"], [TEAL,GOLD,RED]):
    oval(s, nx-0.98, ny-0.56, 1.96, 1.12, c)
    txt(s, lab, nx-0.98, ny-0.56, 1.96, 1.12, 15, WHITE if c!=GOLD else INK, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
def edge(p, q, sh=0.98):
    (ax,ay),(bx,by)=p,q; dx,dy=bx-ax,by-ay; d=math.hypot(dx,dy)
    return ax+dx/d*sh, ay+dy/d*sh, bx-dx/d*sh, by-dy/d*sh
for a,b in [(0,1),(1,2),(2,0)]:
    arrow(s, *edge(pts[a],pts[b]), color=INK, lw=2.5)
txt(s, "learn A → counter with B → revert to A …", 0, 6.45, W, 0.5, 16, INK, italic=True, align=PP_ALIGN.CENTER)
notes(s, "Through human playtesting I found the early bot was greedy — it played point cards instead of "
         "saving them — and exploitable. Worse, always playing itself caused cyclic degradation: it would "
         "learn a strategy, learn a counter, then revert, because that always wins locally.")

# ════════════════════════════════════════════════════════════════════════════
# 9 — A BETTER DOJO  (heuristics + league gate)
# ════════════════════════════════════════════════════════════════════════════
s = slide(dark=True)
kicker(s, "THE FIX: A CURATED OPPONENT POOL")
txt(s, "heuristic bots", 0.95, 1.75, 3.6, 0.4, 14, TEAL, bold=True)
for k,name in enumerate(["aggressive","point-hoarder","random","passive"]):
    yy = 2.3+k*0.72
    rect(s, 0.95, yy, 3.0, 0.56, MID, rounded=True)
    txt(s, name, 0.95, yy, 3.0, 0.56, 14, OFFWHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    arrow(s, 4.0, yy+0.28, 5.35, 4.0, color=PANEL, lw=1.3)
oval(s, 5.35, 3.45, 1.55, 1.1, GOLD)
txt(s, "GATE", 5.35, 3.45, 1.55, 1.1, 16, INK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, "only best /\nmost exploitative", 5.1, 4.65, 2.05, 0.7, 11, MUTED, align=PP_ALIGN.CENTER)
txt(s, "league", 9.4, 1.75, 3.0, 0.4, 14, GOLD, bold=True, align=PP_ALIGN.CENTER)
for k in range(5):
    yy = 4.7-k*0.62; w = 2.0+k*0.25
    rect(s, 10.6-w/2+0.4, yy, w, 0.5, PANEL if k<4 else TEAL, rounded=True)
    txt(s, f"v{k+1}", 10.6-w/2+0.4, yy, w, 0.5, 13, OFFWHT if k<4 else INK, bold=(k==4),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
arrow(s, 6.9, 4.0, 9.3, 3.05, color=GOLD, lw=2)
notes(s, "The fix: heuristic opponents for strategic variety, plus league-based training — only the best "
         "or most exploitative past versions get promoted into the opponent pool. That breaks the cycle "
         "and forces the policy to stay robust.")

# ════════════════════════════════════════════════════════════════════════════
# 10 — WHAT IT SEES  (feature heatmap image)
# ════════════════════════════════════════════════════════════════════════════
s = slide(dark=False)
kicker(s, "WHAT THE POLICY SEES", dark=False)
txt(s, "172", 0, 1.15, W, 1.0, 64, INK, bold=True, align=PP_ALIGN.CENTER, font="Georgia")
txt(s, "features", 0, 2.15, W, 0.4, 16, MUTED, align=PP_ALIGN.CENTER)
pic_center(s, f"{ASSETS}/heat_features.png", cy=4.4, w=10.6)
oval(s, 3.0, 6.45, 0.3, 0.3, GOLD); txt(s, "~60 hand-crafted by me", 3.42, 6.43, 4, 0.35, 15, INK, bold=True)
oval(s, 7.4, 6.45, 0.3, 0.3, TEAL); txt(s, "~112 augmented with Claude", 7.82, 6.43, 5, 0.35, 15, INK, bold=True)
notes(s, "Then: what does the policy actually see? I landed on 172 features — highest card rank, total "
         "points left in the deck, cards remaining, and so on. About 60 were mine; Claude helped augment "
         "the rest. The goal was a signal that's meaningful, directed, and noise-free.")

# ════════════════════════════════════════════════════════════════════════════
# 11 — THE RULES  (card order + bombs)
# ════════════════════════════════════════════════════════════════════════════
s = slide(dark=True)
kicker(s, "THE RULES, IN ONE SLIDE")
txt(s, "weak", 0.95, 1.65, 2, 0.35, 12, MUTED)
txt(s, "strong", 10.4, 1.65, 2.2, 0.35, 12, GOLD, align=PP_ALIGN.RIGHT)
order = [("3","♣",False),("5","♦",True),("10","♥",True),("K","♠",True),("A","♥",False),("2","♣",False)]
ocw=1.18; ostep=1.45; ox=0.95
for k,(rk,su,pt) in enumerate(order):
    card(s, ox+k*ostep, 2.05, ocw, 1.7, rk, su, point=pt)
joker(s, ox+6*ostep, 2.05, ocw, 1.7, red=True)
joker(s, ox+7*ostep, 2.05, ocw, 1.7, red=False)
arrow(s, 1.0, 4.05, 12.3, 4.05, color=PANEL, lw=2)
txt(s, "gold pip = point card  (5 / 10 / K)", 0.95, 4.2, 7, 0.4, 13, GOLD, italic=True)
txt(s, "bombs win the hand anytime:", 0.95, 4.85, 11, 0.4, 14, RED, bold=True)
by=5.3; bcw=0.78; bch=1.1
def mini(x, items):
    for k,(rk,su,pt) in enumerate(items): card(s, x+k*0.5, by, bcw, bch, rk, su, point=pt)
mini(0.95, [("5","♦",True),("10","♥",True),("K","♠",True)])
txt(s, "<", 3.4, by+0.2, 0.6, 0.6, 30, MUTED, align=PP_ALIGN.CENTER, bold=True)
mini(4.1, [("7","♠",False),("7","♥",False),("7","♦",False),("7","♣",False)])
txt(s, "<", 6.95, by+0.2, 0.6, 0.6, 30, MUTED, align=PP_ALIGN.CENTER, bold=True)
joker(s, 7.7, by, bcw, bch, red=True); joker(s, 8.25, by, bcw, bch, red=False)
txt(s, "5-10-K   <   four of a kind   <   two jokers", 9.3, by+0.35, 4, 0.8, 13, OFFWHT, italic=True)
notes(s, "Quick rules for the demo: win point cards — 5s, 10s, Ks. Order runs 3 up to ace, then 2, then "
         "the red and black jokers. Bombs win the hand anytime: 5-10-K, beaten by four of a kind, beaten "
         "by two jokers. Finish first for +20.")

# ════════════════════════════════════════════════════════════════════════════
# 12 — REPORT CARD  (benchmark chart image)
# ════════════════════════════════════════════════════════════════════════════
s = slide(dark=False)
kicker(s, "GOOD, NOT WORLD-CLASS — YET", dark=False)
pic_center(s, f"{ASSETS}/chart_bench.png", cy=4.2, w=9.6)
notes(s, "Performance is good but not great. It beats random about 98% of the time, beats novice and "
         "intermediate humans almost always, and it's about 50-50 against experienced players. Lots of "
         "headroom — and as it improves, even strong players could use it as a training tool.")

# ════════════════════════════════════════════════════════════════════════════
# 13 — DEPLOYMENT  (Balatro-style mockup; swap in real screenshot)
# ════════════════════════════════════════════════════════════════════════════
s = slide(dark=True)
kicker(s, "DEPLOYMENT")
bx,byy,bw,bh = 1.9, 1.9, 9.5, 5.0
rect(s, bx, byy, bw, bh, RGBColor(0x0A,0x0E,0x1A), line=PANEL, lw=1.5, rounded=True)
rect(s, bx, byy, bw, 0.5, MID)
for k,c in enumerate([RED,GOLD,TEAL]): oval(s, bx+0.18+k*0.28, byy+0.15, 0.2, 0.2, c)
rect(s, bx+1.2, byy+0.1, 6.3, 0.3, DARK, rounded=True)
txt(s, "laoban.cards", bx+1.4, byy+0.11, 6, 0.28, 12, MUTED)
oval(s, bx+1.3, byy+0.95, 6.9, 2.5, GREEN, line=RGBColor(0x14,0x55,0x30), lw=2)
card(s, bx+4.2, byy+1.15, 0.8, 1.15, "", "", facedown=True)
card(s, bx+4.9, byy+1.15, 0.8, 1.15, "", "", facedown=True)
for k,(rk,su,pt) in enumerate([("9","♠",False),("J","♦",False),("10","♥",True),("K","♣",True),("2","♠",False)]):
    card(s, bx+3.0+k*0.85, byy+3.5, 0.78, 1.12, rk, su, point=pt)
rect(s, bx+0.4, byy+3.85, 1.5, 0.55, GOLD, rounded=True)
txt(s, "PLAY", bx+0.4, byy+3.85, 1.5, 0.55, 16, INK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
glowbehind(s, "gold", W/2, 7.05, 5)
txt(s, "laoban.cards", 0, 6.75, W, 0.5, 22, GOLD, bold=True, align=PP_ALIGN.CENTER, font="Georgia")
notes(s, "Deployment was the easy part. I loved Balatro as a kid, so I built the UI around that feel and "
         "hosted the backend on Render. NOTE TO SELF: replace this mockup with a real screenshot/clip of "
         "laoban.cards before recording.")

# ════════════════════════════════════════════════════════════════════════════
# 14 — MORE SEATS  (future: heads-up → 4–6 players)
# ════════════════════════════════════════════════════════════════════════════
s = slide(dark=False)
kicker(s, "WHAT'S NEXT", dark=False)
def table(cx, cy, rx, ry, seats, label):
    oval(s, cx-rx, cy-ry, 2*rx, 2*ry, GREEN, line=RGBColor(0x14,0x55,0x30), lw=2)
    for k in range(seats):
        a = -90+k*(360/seats)
        oval(s, cx+(rx+0.45)*math.cos(math.radians(a))-0.32,
                cy+(ry+0.45)*math.sin(math.radians(a))-0.32, 0.64, 0.64, GOLD)
    txt(s, label, cx-2.3, cy+ry+0.6, 4.6, 0.5, 22, INK, bold=True, align=PP_ALIGN.CENTER, font="Georgia")
table(3.4, 3.6, 1.7, 1.1, 2, "now: heads-up")
arrow(s, 6.1, 3.6, 7.5, 3.6, color=GOLD, lw=3)
table(10.0, 3.6, 1.9, 1.25, 6, "next: 4–6 players")
notes(s, "Long-term, I see this as Lichess-for-5/10/K — a low-stakes place to learn, especially in China "
         "where the game is huge. Short-term, I want multiplayer and multi-bot formats. I started heads-up "
         "for a solid MVP, but that's not how it's really played, and strategy shifts a lot with more seats.")

# ════════════════════════════════════════════════════════════════════════════
# 15 — CLOSE
# ════════════════════════════════════════════════════════════════════════════
s = slide(dark=True)
glowbehind(s, "gold", W/2, 2.2, 6)
for k,(rk,su) in enumerate([("5","♦"),("10","♥"),("K","♠")]):
    card(s, 5.2+k*1.5, 1.3, 1.4, 2.0, rk, su, point=True)
txt(s, "let's play.", 0, 3.8, W, 1.0, 60, WHITE, bold=True, align=PP_ALIGN.CENTER, font="Georgia")
txt(s, "laoban.cards", 0, 4.95, W, 0.7, 26, GOLD, bold=True, align=PP_ALIGN.CENTER, font="Georgia")
txt(s, "AI: Claude helped write the encoder, decoder, game environment & training loop, and brainstorm features.",
    1.3, 6.7, 10.7, 0.5, 12, MUTED, italic=True, align=PP_ALIGN.CENTER)
notes(s, "OK — quick demo. Hopefully the bot doesn't embarrass me. Try it yourself at laoban.cards; it's "
         "genuinely addicting. [Demo: play a hand or two, and narrate one strategic decision out loud.]")

prs.save("laoban_slides.pptx")
print("Saved laoban_slides.pptx —", len(prs.slides._sldIdLst), "slides")
